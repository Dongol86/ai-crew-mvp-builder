from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from pptx import Presentation

class DocumentProcessor:
    def process_file(self, file_path):
        if file_path.endswith('.pdf'):
            return self._process_pdf(file_path)
        elif file_path.endswith('.docx'):
            return self._process_docx(file_path)
        elif file_path.endswith('.pptx'):
            return self._process_pptx(file_path)
            
    def _process_pdf(self, file_path):
        loader = PyPDFLoader(file_path)
        return loader.load()
        
    def _process_docx(self, file_path):
        loader = Docx2txtLoader(file_path)
        return loader.load()
        
    def _process_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return text